"""Generate PowerPoint presentation for Spark Cluster Analysis."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

RESULTS_DIR = os.path.dirname(os.path.abspath(__file__))
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Colors
DARK_BLUE = RGBColor(0x1B, 0x2A, 0x4A)
LIGHT_BLUE = RGBColor(0x34, 0x98, 0xDB)
GREEN = RGBColor(0x2E, 0xCC, 0x71)
RED = RGBColor(0xE7, 0x4C, 0x3C)
GRAY = RGBColor(0x7F, 0x8C, 0x8D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def add_title_slide(title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    # Background
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = DARK_BLUE
    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    # Subtitle
    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(20)
    p2.font.color.rgb = LIGHT_BLUE
    p2.alignment = PP_ALIGN.CENTER


def add_content_slide(title, bullets, image_path=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Title bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(12), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    if image_path and os.path.exists(os.path.join(RESULTS_DIR, image_path)):
        # Image on right, bullets on left
        left_width = Inches(5.5) if bullets else Inches(12)
        if bullets:
            txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(5.5), Inches(5.5))
            tf2 = txBox2.text_frame
            tf2.word_wrap = True
            for bullet in bullets:
                p = tf2.add_paragraph()
                p.text = bullet
                p.font.size = Pt(16)
                p.space_after = Pt(8)
            slide.shapes.add_picture(
                os.path.join(RESULTS_DIR, image_path),
                Inches(6.3), Inches(1.2), Inches(6.5), Inches(5.5))
        else:
            slide.shapes.add_picture(
                os.path.join(RESULTS_DIR, image_path),
                Inches(1), Inches(1.3), Inches(11), Inches(5.8))
    elif bullets:
        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12), Inches(5.5))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        for bullet in bullets:
            p = tf2.add_paragraph()
            p.text = bullet
            p.font.size = Pt(18)
            p.space_after = Pt(10)


def add_table_slide(title, headers, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Title bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(12), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Table
    num_rows = len(rows) + 1
    num_cols = len(headers)
    left = Inches(0.5)
    top = Inches(1.3)
    width = Inches(12.3)
    height = Inches(0.4) * num_rows

    table = slide.shapes.add_table(num_rows, num_cols, left, top, width, height).table

    # Header
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.color.rgb = WHITE
            paragraph.font.bold = True
            paragraph.font.size = Pt(12)

    # Rows
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)


# === BUILD PRESENTATION ===

# Slide 1: Title
add_title_slide(
    "Multi-Model Distributed Inference",
    "Spark Cluster Benchmark — POC Results\nJuly 2026 | AIA Team"
)

# Slide 2: What We Built
add_content_slide("What We Built", [
    "10 ML models (signal classification, image recognition, object detection)",
    "3 inference modes: Single GPU, Hybrid CPU+GPU, Distributed Spark",
    "AWS cluster: m5.xlarge (master) + g4dn.xlarge (GPU worker)",
    "25 benchmark runs across GPU/CPU/Hybrid with increasing data",
    "Automated deployment via CDK + deploy script",
    "Complete monitoring: CloudWatch + Spark UI + REST API capture",
])

# Slide 3: Architecture
add_content_slide("Cluster Architecture", [
    "Master (10.0.0.187): Spark Master + Driver + CPU Worker",
    "GPU Worker (10.0.0.45): Spark Worker + Tesla T4 GPU",
    "Communication: Spark RPC over private VPC (port 7077)",
    "Data distribution: RDD partitioning (not broadcast)",
    "Model weights: Broadcast once (~75MB), cached on executors",
    "mapPartitions: Models loaded ONCE per executor, reused across tasks",
])

# Slide 4: Mode Comparison Chart
add_content_slide("Mode Throughput Comparison", None, "cluster_mode_comparison.png")

# Slide 5: Results Table
add_table_slide("Benchmark Results — All Modes at Maximum Load",
    ["Mode", "Samples", "Partitions", "Throughput", "Time", "GPU Speedup"],
    [
        ["gpu_only", "61,700", "6", "2,923 /sec", "21.1s", "1.8x"],
        ["cpu_only", "61,700", "6", "1,588 /sec", "38.9s", "baseline"],
        ["hybrid", "61,700", "6", "2,120 /sec", "29.1s", "1.3x"],
        ["Single GPU*", "25,700", "—", "30,087 /sec", "0.85s", "22x"],
        ["Hybrid GPU*", "25,700", "—", "29,910 /sec", "0.86s", "22x"],
    ])

# Slide 6: Scaling Chart
add_content_slide("Throughput Scaling with Data Size", [
    "GPU advantage grows with more data",
    "5K samples: GPU slower (model load overhead)",
    "15K: breakeven point",
    "25K+: GPU 1.5-1.8x faster than CPU",
    "Bottleneck: CPU executor (9.14s vs GPU 0.33s)",
], "cluster_scaling_by_mode.png")

# Slide 7: Executor Distribution
add_content_slide("Task Distribution Across Executors", [
    "GPU executor: 5 of 6 tasks (83%)",
    "CPU executor: 1 of 6 tasks (17%)",
    "Spark auto-routes more work to faster executor",
    "GPU per-task: 0.32s (32,128 samples/sec)",
    "CPU per-task: 9.14s (1,125 samples/sec)",
    "GPU is 13-29x faster per task",
], "cluster_executor_distribution.png")

# Slide 8: Task Timeline
add_content_slide("Task Execution Timeline (6 Partitions)", [
    "GPU Executor: P0(0.69s) → P2(0.32s) → P3(0.33s) → P4(0.33s) → P5(0.34s)",
    "CPU Executor: P1 ─────────── 9.14s ───────────",
    "",
    "GPU total: 2.0s for 51,419 samples",
    "CPU total: 9.14s for 10,281 samples",
    "Wall clock = max(GPU, CPU) + overhead = 21.1s",
    "",
    "INSIGHT: Total throughput limited by slowest executor",
], "final_stage_timeline.png")

# Slide 9: mapPartitions Optimization
add_content_slide("mapPartitions: Load Models Once", [
    "OLD (map): Each task loads 10 models → 1.6s × 6 tasks = 9.6s wasted",
    "NEW (mapPartitions): Load once, reuse for all tasks on that executor",
    "",
    "Partition 0 (first): Model load = 1.60s (cold)",
    "Partition 2 (second): Model load = 0.53s (cached) → 3.1x faster",
    "Partition 3-5: Model load = 0.51s (reused)",
    "",
    "Total savings: ~5s on GPU executor",
    "This is why distributed throughput improved from 1,880 → 2,923 /sec",
])

# Slide 10: Model Load vs Inference
add_content_slide("Model Load vs Inference Time", None, "cluster_load_vs_inference.png")

# Slide 11: GPU Confirmation
add_content_slide("GPU Confirmed in Spark Executor", [
    "From executor stderr logs on 10.0.0.45:",
    "",
    '  [Executor] partition=0, host=ip-10-0-0-45, cuda=True, device=cuda',
    "",
    "Proof:",
    "  - torch.cuda.is_available() = True inside executor subprocess",
    "  - Models loaded on device='cuda' (Tesla T4)",
    "  - Inference runs on GPU (0.32s vs 9.14s on CPU)",
    "",
    "Fix applied: NVIDIA_VISIBLE_DEVICES=all + LD_LIBRARY_PATH in executor env",
])

# Slide 12: Per-Partition Throughput
add_content_slide("Per-Partition Throughput (GPU vs CPU)", None, "cluster_partition_throughput.png")

# Slide 13: Challenges
add_content_slide("Key Challenges Solved", [
    "1. GPU not visible in Spark executor → Set NVIDIA env vars in Spark config",
    "2. Driver OOM on large data → Increase partitions, use mapPartitions",
    "3. Task serialization > 512MB → Increase spark.rpc.message.maxSize",
    "4. Per-task model loading overhead → mapPartitions (load once)",
    "5. CPU bottleneck limits cluster throughput → Need all-GPU cluster",
    "6. SSM agent crashes under heavy load → Use larger instances",
    "7. Docker GPU passthrough to subprocess → Fixed with env vars",
])

# Slide 14: Air-gapped Projection
add_content_slide("Air-Gapped 5-Node Production Projection", [
    "Hardware: 5 nodes × 256GB RAM × 24GB VRAM × 4TB disk",
    "",
    "Single GPU (1 node): ~40,000-50,000 samples/sec",
    "Distributed 5 GPUs (current code): ~50,000-65,000 samples/sec",
    "Distributed optimized (persistent): ~150,000+ samples/sec",
    "",
    "No CPU bottleneck (all nodes have GPU)",
    "No Docker GPU issues (native Spark install)",
    "No memory limits (256GB RAM per node)",
    "Standalone mode — simple, reliable, no extra infrastructure",
], "final_projected_airgapped.png")

# Slide 15: Recommendations
add_content_slide("Recommendations for Production", [
    "1. Use Standalone mode (not YARN) — simplest for dedicated 5-node cluster",
    "2. Install Spark natively (no Docker) — eliminates GPU subprocess issues",
    "3. Set partitions = 5 (1 per GPU node) — each loads models once",
    "4. Use mapPartitions — models persist across batches",
    "5. Pre-bake model weights in deployment package — no downloads",
    "6. Set spark.rpc.message.maxSize=2048 — large data handled easily",
    "7. Enable NVIDIA MPS for multi-process GPU sharing",
    "8. Monitor via Spark UI + nvidia-smi scripts (no CloudWatch in air-gap)",
])

# Slide 16: Summary
add_title_slide(
    "POC Complete",
    "GPU distributed inference proven\n"
    "2,923 samples/sec (2-node) → projected 50,000+ (5-node)\n"
    "Zero task failures across 25 runs\n"
    "Framework ready for AIA team deployment"
)

# Save
output_path = os.path.join(RESULTS_DIR, "Cluster_Analysis_Presentation.pptx")
prs.save(output_path)
print(f"Presentation saved: {output_path}")
